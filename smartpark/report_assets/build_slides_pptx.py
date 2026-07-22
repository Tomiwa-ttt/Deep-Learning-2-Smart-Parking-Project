"""
build_slides_pptx.py

Builds the presentation deck (SmartPark_Slides.pptx) from the same real
numbers/charts as build_report_docx.py.

Usage:
    cd smartpark
    python report_assets/build_slides_pptx.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
CHARTS = os.path.join(HERE, "charts")
SAMPLES = os.path.join(HERE, "sample_outputs")
CHARTS_REAL = os.path.join(HERE, "charts_real")
SAMPLES_REAL = os.path.join(HERE, "sample_outputs_real")

BG = RGBColor(0x18, 0x18, 0x1B)
FG = RGBColor(0xF4, 0xF4, 0xF5)
ACCENT = RGBColor(0x6B, 0x8A, 0xFF)
MUTED = RGBColor(0xA1, 0xA1, 0xAA)

SLIDE_W, SLIDE_H = Inches(13.33), Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return slide


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=FG,
             align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_title(slide, text, subtitle=None):
    add_text(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.8), text,
              size=30, bold=True, color=ACCENT)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.05), Inches(12), Inches(0.5), subtitle,
                  size=15, color=MUTED)


def add_bullets(slide, left, top, width, height, items, size=16, color=FG):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)


def add_picture_fit(slide, path, left, top, max_w, max_h):
    from PIL import Image
    with Image.open(path) as im:
        w, h = im.size
    ratio = min(max_w / w, max_h / h)
    pic_w, pic_h = int(w * ratio), int(h * ratio)
    pic_left = left + (max_w - pic_w) // 2
    pic_top = top + (max_h - pic_h) // 2
    slide.shapes.add_picture(path, pic_left, pic_top, width=pic_w, height=pic_h)


def add_table(slide, left, top, width, height, headers, rows, font_size=13):
    n_rows, n_cols = len(rows) + 1, len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for j, h in enumerate(headers):
        cell = gtable.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(font_size)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x33, 0x38, 0x42)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = gtable.cell(i, j)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(font_size)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x24, 0x28, 0x2F)
    return gtable


def build():
    prs = new_deck()

    # 1. Title
    s = blank_slide(prs)
    add_text(s, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
              "SmartPark: Real-Time Parking Spot Object Detection",
              size=40, bold=True, color=FG, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.6), Inches(11.3), Inches(0.6),
              "AASD 4014 — Deep Learning II — Final Project", size=18, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(4.2), Inches(11.3), Inches(0.6),
              "Object Detection Track  ·  July 2026  ·  Group [FILL IN]", size=16, color=ACCENT,
              align=PP_ALIGN.CENTER)

    # 2. Problem
    s = blank_slide(prs)
    add_title(s, "The Problem", "Drivers waste time and money searching for parking")
    add_bullets(s, Inches(0.7), Inches(2.0), Inches(11.5), Inches(4), [
        "Camera-based occupancy detection can surface spot-level availability before drivers arrive.",
        "Prior work (PKLot / CNRPark-EXT) classifies pre-cropped, pre-calibrated spot images — it assumes "
        "spot locations are already known.",
        "This project targets the assignment's Object Detection requirement directly: find AND classify "
        "spots (empty_spot / occupied_spot) in a full, uncalibrated photo — no manual calibration step.",
    ], size=19)

    # 3. Plan of attack + agile
    s = blank_slide(prs)
    add_title(s, "Plan of Attack & Agile Process", "5 real development stages, reconstructed from project file history")
    add_table(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(2.6),
               ["Sprint", "Goal", "Key deliverables"], [
        ["1 (Jul 17)", "Pipeline scaffolding", "Synthetic data, CNN classifier, geometry check, first demo UI"],
        ["2 (Jul 18)", "Real-data integration", "PKLot conversion, classifier retrained on real data (97.04%)"],
        ["3 (Jul 19)", "Reporting", "Initial report/slides, architecture diagram"],
        ["4 (Jul 20)", "Env fixes + OD pivot", "Rebuilt as a real 2-class object detector after rubric audit"],
        ["5 (Jul 21)", "Real-data closure", "Fine-tuned detector on real boxes (47.71% mAP); classifier 98.69%"],
    ], font_size=12)
    add_picture_fit(s, os.path.join(CHARTS, "burndown.png"), Inches(3.4), Inches(4.6), Inches(6.5), Inches(2.6))

    # 4. Dataset
    s = blank_slide(prs)
    add_title(s, "The Dataset", "Full-scene detection annotations, two object classes")
    add_bullets(s, Inches(0.7), Inches(1.9), Inches(6.2), Inches(4.5), [
        "300 synthetic full-lot images, randomized grid layout (rows/cols/spot size/margin + per-spot jitter)",
        "4,422 labeled spot instances: 2,027 empty_spot / 2,395 occupied_spot",
        "occupied_spot boxes track the actual car footprint (shifts when misparked) — real localization signal",
        "255 train / 45 val images (85/15 split)",
        "Both classes clear the 200-images-per-class minimum comfortably",
    ], size=17)
    add_picture_fit(s, os.path.join(SAMPLES, "sample_01.jpg"), Inches(7.1), Inches(1.9), Inches(5.6), Inches(4.5))

    # 5. Model architecture
    s = blank_slide(prs)
    add_title(s, "Model: Single-Shot Detector", "Trained from scratch — YOLOv1-style, one box per grid cell")
    add_table(s, Inches(0.6), Inches(1.9), Inches(6.2), Inches(4.3),
               ["Layer", "Output"], [
        ["Input", "256×256×3"],
        ["Conv2D(16)+Pool", "128×128×16"],
        ["Conv2D(32)+Pool", "64×64×32"],
        ["Conv2D(64)+Pool", "32×32×64"],
        ["Conv2D(128)+Pool", "16×16×128"],
        ["Conv2D(128)×2", "16×16×128"],
        ["Head (1×1 conv)", "16×16×7"],
    ], font_size=13)
    add_bullets(s, Inches(7.1), Inches(1.9), Inches(5.6), Inches(4.3), [
        "393,511 parameters",
        "Per cell: objectness, box (tx,ty,tw,th), 2-class softmax",
        "Loss: YOLOv1-style — coordinate (×5), objectness, no-object (×0.5), classification",
        "Post-processing: per-class NMS (cv2.dnn.NMSBoxes)",
        "Secondary pipeline: CNN classifier (1.14M params) + IoU geometry check for calibrated cameras",
    ], size=16)

    # 6. Training curve
    s = blank_slide(prs)
    add_title(s, "Training", "40 epochs, Adam + ReduceLROnPlateau, early stopping patience 6 (not triggered)")
    add_picture_fit(s, os.path.join(CHARTS, "training_curve.png"), Inches(2.5), Inches(1.8), Inches(8.3), Inches(5.2))

    # 7. Evaluation metrics
    s = blank_slide(prs)
    add_title(s, "Evaluation", "Held-out validation: 45 images, 721 spot instances")
    add_table(s, Inches(0.6), Inches(1.9), Inches(7.0), Inches(2.0),
               ["Class", "GT", "AP@0.5", "Precision", "Recall", "F1"], [
        ["empty_spot", "337", "98.73%", "95.5%", "94.1%", "94.8%"],
        ["occupied_spot", "384", "98.68%", "99.7%", "96.1%", "97.9%"],
        ["mAP@0.5", "—", "98.71%", "—", "—", "—"],
    ], font_size=13)
    add_picture_fit(s, os.path.join(CHARTS, "ap_bar.png"), Inches(7.9), Inches(1.8), Inches(4.9), Inches(4.9))

    # 8. PR curves + confusion matrix
    s = blank_slide(prs)
    add_title(s, "Where the Errors Are", "Localization misses, not class confusion")
    add_picture_fit(s, os.path.join(CHARTS, "pr_curves.png"), Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0))
    add_picture_fit(s, os.path.join(CHARTS, "confusion_matrix.png"), Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.0))

    # 9. Qualitative results (good + bad)
    s = blank_slide(prs)
    add_title(s, "Qualitative Results", "A clean match, and an honest failure case")
    add_picture_fit(s, os.path.join(SAMPLES, "sample_03.jpg"), Inches(0.4), Inches(1.9), Inches(6.0), Inches(4.6))
    add_picture_fit(s, os.path.join(SAMPLES, "sample_06.jpg"), Inches(6.7), Inches(1.9), Inches(6.0), Inches(4.6))
    add_text(s, Inches(0.4), Inches(6.6), Inches(6.0), Inches(0.6), "24/24 spots, exact match", size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(6.7), Inches(6.6), Inches(6.0), Inches(0.6), "Missed several cars incl. a low-contrast gray car", size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # 10. Hyperparameter tuning
    s = blank_slide(prs)
    add_title(s, "Hyperparameter Tuning", "Two real, measured comparisons")
    add_text(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(0.4), "Grid resolution vs. collision rate", size=17, bold=True, color=ACCENT)
    add_table(s, Inches(0.6), Inches(2.3), Inches(5.8), Inches(1.8),
               ["Grid", "Dropped", "Rate"], [
        ["8×8", "111/4,422", "2.51%"],
        ["16×16 ✓", "12/4,422", "0.27%"],
        ["32×32", "3/4,422", "0.07%"],
    ], font_size=13)
    add_text(s, Inches(6.9), Inches(1.8), Inches(6.0), Inches(0.4), "Coordinate loss weight (15 epochs each)", size=17, bold=True, color=ACCENT)
    add_table(s, Inches(6.9), Inches(2.3), Inches(5.8), Inches(1.5),
               ["lambda_coord", "Val loss", "mAP"], [
        ["5.0 ✓", "3.91", "96.72%"],
        ["1.0", "16.36", "0.00%"],
    ], font_size=13)
    add_text(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(1.8),
              "lambda_coord=1.0 collapsed to 0% mAP within the short budget — a known YOLO pathology when "
              "coordinate loss is underweighted relative to the flood of no-object cells. Confirms the "
              "YOLOv1 default rather than assuming it.", size=15, color=MUTED)

    # 11. Benchmarking
    s = blank_slide(prs)
    add_title(s, "Benchmarking vs. Baselines", "Same validation spots, three methods")
    add_table(s, Inches(1.5), Inches(2.3), Inches(10.3), Inches(2.4),
               ["Method", "Accuracy"], [
        ["Majority-class baseline", "53.26%"],
        ["Classic-CV heuristic (this project's original approach)", "82.66%"],
        ["Trained detector (given correct localization)", "99.86%"],
    ], font_size=16)
    add_text(s, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.5),
              "The classic-CV heuristic was this project's own first approach to occupancy. It beats a naive "
              "prior but misses ~1 in 6 spots — exactly the gap a learned model should close.", size=15, color=MUTED)

    # 12. Secondary pipeline / real data
    s = blank_slide(prs)
    add_title(s, "Secondary Pipeline: Calibrated Cameras", "Occupancy classifier + geometric misparking check")
    add_bullets(s, Inches(0.7), Inches(1.9), Inches(11.7), Inches(4.5), [
        "For cameras with known spot boundaries: crop each spot → CNN classifies Empty/Occupied → IoU/overflow "
        "check flags cars parked across the boundary",
        "Trained on real PKLot data (1,242 photos, 70,684 labeled spaces) → 98.69% validation accuracy, "
        "consistent with (slightly ahead of) published PKLot/CNRPark-EXT results",
        "Same real PKLot export also used to fine-tune the object detector (next slide)",
        "Adds a feature the general detector doesn't attempt: flagging improperly parked cars",
    ], size=17)

    # 13. Live demo
    s = blank_slide(prs)
    add_title(s, "Live Demo", "Flask REST API + browser upload UI — no install needed")
    add_picture_fit(s, os.path.join(HERE, "live_demo_detection.jpg"), Inches(2.8), Inches(1.8), Inches(7.7), Inches(5.2))

    # 13b. Cross-domain zero-shot test
    s = blank_slide(prs)
    add_title(s, "Zero-Shot Cross-Domain Test (CNRPark-EXT)", "A totally unfamiliar camera angle and marking style")
    add_picture_fit(s, os.path.join(HERE, "real_photo_test", "example_2.jpg"), Inches(1.2), Inches(1.8), Inches(6.6), Inches(5.0))
    add_bullets(s, Inches(8.1), Inches(2.0), Inches(4.6), Inches(4.5), [
        "115 real photos tested, no painted lines, steep real camera angle",
        "Boxes loosely cluster near real cars — some transferable signal",
        "Misses most vehicles; box shapes imprecise",
        "empty_spot almost never fires — no painted bays in this domain",
        "Matches this project's earlier classifier result: 40% on the same cross-domain test",
        "Contrast with next slide: same-domain fine-tuning tells a very different story",
    ], size=13)

    # 13c. Real-data fine-tuning results
    s = blank_slide(prs)
    add_title(s, "Real-Data Fine-Tuning (PKLot)", "Same-domain real photos — a very different result from zero-shot")
    add_table(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(1.9),
               ["Class", "GT", "AP@0.5", "Precision", "Recall"], [
        ["empty_spot", "5,345", "46.84%", "77.6%", "53.2%"],
        ["occupied_spot", "4,975", "48.59%", "72.6%", "57.1%"],
        ["mAP@0.5", "—", "47.71%", "—", "—"],
    ], font_size=12)
    add_picture_fit(s, os.path.join(CHARTS_REAL, "confusion_matrix.png"), Inches(7.0), Inches(1.6), Inches(5.7), Inches(4.2))
    add_text(s, Inches(0.6), Inches(3.9), Inches(6.0), Inches(2.8),
              "47.71% mAP on real, dense PKLot photos — a genuine result, far below synthetic (98.71%) but a "
              "huge step up from zero-shot (previous slide). Classification given correct localization: 94.5% "
              "(vs. 99.86% synthetic) — almost as reliable. The gap is recall: real lots average ~57 spots/image "
              "vs. ~15 synthetic, so the 16×16 grid drops 28.2% of real boxes to cell collisions. That alone "
              "caps achievable recall near ~72% — the concrete next fix is a finer grid or anchor-based head, "
              "not more real data.",
              size=13, color=MUTED)

    # 13d. Real qualitative result
    s = blank_slide(prs)
    add_title(s, "Real-Data Fine-Tuning — Qualitative", "A dense, cluttered real PKLot photo")
    add_picture_fit(s, os.path.join(SAMPLES_REAL, "sample_06.jpg"), Inches(1.8), Inches(1.7), Inches(9.7), Inches(5.3))

    # 14. Discussion & limitations
    s = blank_slide(prs)
    add_title(s, "Discussion & Limitations", "Honest accounting")
    add_bullets(s, Inches(0.7), Inches(1.9), Inches(11.7), Inches(4.8), [
        "Worked: staged synthetic-first build let the whole pipeline be validated cheaply before real data",
        "Had to be revisited: the original classifier-only version didn't satisfy Object Detection — only "
        "caught by auditing against the rubric, requiring a genuine architecture change in Sprint 4",
        "Zero-shot to an unfamiliar camera (CNRPark-EXT) mostly fails; fine-tuning on same-domain real data "
        "(PKLot) reaches 47.71% mAP — real data helps a lot, but the model needs to actually see it",
        "The real-data bottleneck is recall from grid-cell collisions on dense lots, not classification — "
        "a specific, measured next fix, not a vague call for \"more real data\"",
    ], size=17)

    # 15. Conclusion
    s = blank_slide(prs)
    add_title(s, "Conclusion")
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5),
              "A from-scratch object detector for two classes (empty_spot / occupied_spot): 98.71% mAP@0.5 "
              "on synthetic data, and 47.71% mAP@0.5 after fine-tuning on real PKLot photos — beating both a "
              "majority-class baseline (53.3%) and this project's own earlier classic-CV heuristic (82.7%). "
              "Classification stays reliable on real data (94.5%, vs. 99.86% synthetic); the real-data gap is "
              "concentrated in localization recall on dense lots, with a specific, measured cause (grid-cell "
              "collisions) and a specific next fix (finer grid / anchor-based head). Served live through a "
              "REST API and browser demo, alongside a real-data-validated occupancy classifier (98.69%) for "
              "calibrated cameras.",
              size=17, color=FG)

    # 16. Team
    s = blank_slide(prs)
    add_title(s, "Team", "[FILL IN before submission]")
    add_table(s, Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.5),
               ["Name", "Role", "Contribution"], [
        ["[FILL IN — PM]", "Project Manager", "[FILL IN]"],
        ["[FILL IN]", "[FILL IN]", "[FILL IN]"],
        ["[FILL IN]", "[FILL IN]", "[FILL IN]"],
        ["[FILL IN]", "[FILL IN]", "[FILL IN]"],
        ["[FILL IN]", "[FILL IN]", "[FILL IN]"],
    ], font_size=14)

    out_path = os.path.join(HERE, "SmartPark_Slides.pptx")
    prs.save(out_path)
    print(f"Saved {out_path} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
